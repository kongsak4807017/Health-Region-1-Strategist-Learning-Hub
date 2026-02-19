import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

TITLE_MAP = {
    'Title': 'ชื่อเรื่อง',
    'Learning Objectives': 'วัตถุประสงค์การเรียนรู้',
    'Agenda': 'Agenda (หัวข้อวันนี้)',
    'Wrap-up': 'สรุป',
    'Q&A': 'ถาม-ตอบ (Q&A)',
}


def normalize_title(title: str) -> str:
    return TITLE_MAP.get(title, title)


def parse_decks(text):
    lines = text.splitlines()
    decks = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith('### ') and 'Slide Deck' in line:
            # Extract deck id like "4.1" or "5.1"
            m = re.match(r'###\s+(\d+\.\d+)\s+Slide Deck', line)
            deck_id = m.group(1) if m else None
            deck = {'id': deck_id, 'title': line, 'slides': []}
            i += 1
            while i < len(lines):
                line2 = lines[i].strip()
                if line2.startswith('### ') and 'Slide Deck' in line2:
                    i -= 1
                    break
                if line2.startswith('Slide '):
                    # Slide title after em dash or hyphen
                    if '—' in line2:
                        title = line2.split('—', 1)[1].strip()
                    elif '-' in line2:
                        title = line2.split('-', 1)[1].strip()
                    else:
                        title = line2
                    title = normalize_title(title)
                    goal = ''
                    keypoints = ''
                    visual = ''
                    j = i + 1
                    while j < len(lines):
                        l3 = lines[j].strip()
                        if l3.startswith('Slide ') or (l3.startswith('### ') and 'Slide Deck' in l3):
                            j -= 1
                            break
                        if l3.startswith('Goal:'):
                            goal = l3.split(':', 1)[1].strip()
                        elif l3.startswith('Key points:'):
                            keypoints = l3.split(':', 1)[1].strip()
                        elif l3.startswith('Visual:'):
                            visual = l3.split(':', 1)[1].strip()
                        j += 1
                    deck['slides'].append({'title': title, 'goal': goal, 'keypoints': keypoints, 'visual': visual})
                    i = j
                i += 1
            decks.append(deck)
        i += 1
    return decks


def add_slide(prs, title, goal, keypoints, visual):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    def add_bullet(text):
        p = body.add_paragraph()
        p.text = text
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.name = 'Calibri'

    # First paragraph should be the first bullet, not empty
    first = body.paragraphs[0]
    first.text = f"เป้าหมาย: {goal}" if goal else "เป้าหมาย:"
    for run in first.runs:
        run.font.size = Pt(24)
        run.font.name = 'Calibri'

    add_bullet(f"ประเด็นหลัก: {keypoints}" if keypoints else "ประเด็นหลัก:")
    add_bullet(f"ภาพ/สื่อ: {visual}" if visual else "ภาพ/สื่อ:")


def build_pptx(decks, out_dir, file_map):
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    for deck in decks:
        deck_id = deck['id']
        if deck_id not in file_map:
            continue
        prs = Presentation()
        # Set default font sizes by adjusting slide master? We'll keep per-slide.
        for slide in deck['slides']:
            add_slide(prs, slide['title'], slide['goal'], slide['keypoints'], slide['visual'])
        out_path = out_dir / file_map[deck_id]
        prs.save(out_path)


def read_text_with_fallback(path: Path) -> str:
    try:
        return path.read_text(encoding='utf-8')
    except UnicodeDecodeError:
        return path.read_text(encoding='cp874')


def main():
    # Module 1
    m1_text = read_text_with_fallback(Path('Module-1/Module-1-Content.md'))
    m1_decks = parse_decks(m1_text)
    m1_map = {
        '1.1': 'M1-01-Data-Sources.pptx',
        '1.2': 'M1-02-Data-Collection.pptx',
        '1.3': 'M1-03-Data-Verification.pptx',
    }
    build_pptx(m1_decks, 'Module-1', m1_map)

    # Module 2
    m2_text = read_text_with_fallback(Path('Module-2/Module-2-Content.md'))
    m2_decks = parse_decks(m2_text)
    m2_map = {
        '2.1': 'M2-01-RCA-Concepts.pptx',
        '2.2': 'M2-02-RCA-Tools.pptx',
        '2.3': 'M2-03-Risk-Factors.pptx',
    }
    build_pptx(m2_decks, 'Module-2', m2_map)

    # Module 3
    m3_text = read_text_with_fallback(Path('Module-3/Module-3-Content.md'))
    m3_decks = parse_decks(m3_text)
    m3_map = {
        '3.1': 'M3-01-Policy-Frameworks.pptx',
        '3.2': 'M3-02-Stakeholder-Analysis.pptx',
        '3.3': 'M3-03-Policy-Options-MCDA.pptx',
    }
    build_pptx(m3_decks, 'Module-3', m3_map)

    # Module 4
    m4_text = read_text_with_fallback(Path('Module-4/Module-4-Content.md'))
    m4_decks = parse_decks(m4_text)
    m4_map = {
        '4.1': 'M4-01-Recommendation-Structure.pptx',
        '4.2': 'M4-02-Effective-Writing.pptx',
        '4.3': 'M4-03-Policy-Communication.pptx',
    }
    build_pptx(m4_decks, 'Module-4', m4_map)

    # Module 5
    m5_text = read_text_with_fallback(Path('Module-5/Module-5-Content.md'))
    m5_decks = parse_decks(m5_text)
    m5_map = {
        '5.1': 'M5-01-Capstone-Project.pptx',
    }
    build_pptx(m5_decks, 'Module-5', m5_map)


if __name__ == '__main__':
    main()
